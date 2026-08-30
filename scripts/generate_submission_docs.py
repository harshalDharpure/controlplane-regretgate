"""Generate Round 2 submission PDFs and PPTX into submission/."""

from __future__ import annotations

from pathlib import Path

from fpdf import FPDF
from fpdf.enums import XPos, YPos
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission"
OUT.mkdir(exist_ok=True)

INK = (18, 20, 26)
MUTED = (90, 99, 110)
ACCENT = (26, 77, 109)


class Doc(FPDF):
    def __init__(self, subtitle: str):
        super().__init__(format="A4")
        self.subtitle = subtitle
        self.set_auto_page_break(True, 16)
        self.set_margins(16, 16, 16)
        self.add_page()

    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica", "", 8)
        self.set_text_color(*MUTED)
        self.cell(0, 8, f"{self.subtitle}  |  Page {self.page_no()}", align="C")

    def title_block(self, title: str, line: str = ""):
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "B", 20)
        self.set_text_color(*INK)
        self.multi_cell(0, 10, title, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        if line:
            self.set_font("Helvetica", "I", 10)
            self.set_text_color(*MUTED)
            self.multi_cell(0, 6, line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.ln(3)
        y = self.get_y()
        self.set_draw_color(180, 186, 194)
        self.line(self.l_margin, y, self.w - self.r_margin, y)
        self.ln(4)

    def section(self, text: str, size: int = 13):
        self.ln(2)
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "B", size)
        self.set_text_color(*INK)
        self.multi_cell(0, 7, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.ln(1)

    def p(self, text: str):
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*INK)
        self.multi_cell(0, 5.4, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.ln(1.2)

    def b(self, text: str):
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*INK)
        self.multi_cell(0, 5.4, f"  -  {text}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    def code(self, text: str):
        self.set_x(self.l_margin)
        self.set_font("Courier", "", 9)
        self.set_fill_color(244, 245, 247)
        self.set_text_color(*INK)
        self.multi_cell(0, 5, text, fill=True, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.ln(1.5)

    def kv(self, rows: list[tuple[str, str]]):
        for k, v in rows:
            self.set_x(self.l_margin)
            self.set_font("Helvetica", "B", 9)
            self.set_text_color(*ACCENT)
            self.cell(48, 5.5, k)
            self.set_font("Helvetica", "", 9)
            self.set_text_color(*INK)
            self.multi_cell(0, 5.5, v, new_x=XPos.LMARGIN, new_y=YPos.NEXT)


def build_readme_pdf(path: Path) -> None:
    pdf = Doc("RegretGate README")
    pdf.title_block(
        "RegretGate - ControlPlane Checker",
        "Accenture Innovation Challenge Round 2  |  Problem Track 1: ControlPlane.ai",
    )
    pdf.p("Act with confidence. Verify only when the regret is high.")
    pdf.p("Public GitHub: https://github.com/harshalDharpure/controlplane-regretgate")

    pdf.section("1. What this project is about")
    pdf.p(
        "RegretGate is an action-aware AI control plane. Enterprises use generative AI "
        "across chatbots, internal copilots, and regulated decision tools. Failures often "
        "appear only AFTER the model has already acted - wrong answers, wasted compute, "
        "or privacy/policy breaches."
    )
    pdf.p(
        "RegretGate sits BEFORE commit (before reply / send / refund / approve / execute). "
        "For every pending AI action it: identifies intent, estimates Expected Regret, "
        "runs a responsibility check, routes through an intervention ladder, logs an audit "
        "receipt, and can escalate to a human."
    )
    pdf.code("Expected Regret = P(failure) x Impact x Irreversibility   (score 0-100)")

    pdf.section("2. Three silent failure modes")
    pdf.b("Confidently wrong - hallucinations / poor decisions with certainty (grounding signals)")
    pdf.b("Quietly expensive - excess tokens, tool calls, retries, agent thrash")
    pdf.b("Subtly unsafe - bias, leakage, policy violations (PII, secrets, unsafe content)")

    pdf.section("3. How it is helpful")
    pdf.b("Business owners: irreversible actions held for proof or human review")
    pdf.b("Risk / privacy / compliance: hard blocks, geo policy packs, full audit trail")
    pdf.b("AI platform teams: one gate across use cases with different latency/risk appetites")
    pdf.b("FinOps: thrash detection stops quietly expensive agent loops")
    pdf.b("Human reviewers: HITL only when regret is high - less alert fatigue")
    pdf.p(
        "In short: faster where it is safe, stricter where consequences matter - "
        "shifting AI safety from post-hoc postmortems to pre-commit control."
    )

    pdf.section("4. Prototype surfaces")
    pdf.kv(
        [
            ("Control Plane /", "Paste or load pending action; score, ladder, rewrite, receipts"),
            ("Scenarios /scenarios", "One-click enterprise demos (support, copilot, decision)"),
            ("Ops /ops", "HITL queue, audit log, metrics, feedback, policy tuning"),
        ]
    )
    pdf.p("Tech: Next.js 16, React 19, TypeScript, Tailwind, Vitest. No API keys required.")

    pdf.section("5. Prerequisites")
    pdf.b("Node.js 20 or newer (https://nodejs.org)")
    pdf.b("npm (included with Node.js)")
    pdf.b("Git (to clone the repository)")
    pdf.code("node -v\nnpm -v")

    pdf.section("6. How to run the complete application")
    pdf.p("Step 1 - Clone")
    pdf.code(
        "git clone https://github.com/harshalDharpure/controlplane-regretgate.git\n"
        "cd controlplane-regretgate"
    )
    pdf.p("Step 2 - Install dependencies")
    pdf.code("npm install")
    pdf.p("Step 3 - Start the development server")
    pdf.code("npm run dev")
    pdf.p("Step 4 - Open in browser")
    pdf.b("Control Plane: http://localhost:3000")
    pdf.b("Scenarios:     http://localhost:3000/scenarios")
    pdf.b("Ops:           http://localhost:3000/ops")
    pdf.p("Step 5 - Stop server with Ctrl+C in the terminal.")

    pdf.section("7. Optional commands")
    pdf.code("npm run build\nnpm start\nnpm test\nnpm run lint")

    pdf.section("8. Suggested demo walkthrough (3-5 minutes)")
    pdf.b("Control Plane -> Load Safe FAQ reply -> Evaluate -> Pass instantly")
    pdf.b("Scenarios -> Quietly expensive tool thrash -> thrash + elevated regret")
    pdf.b("Scenarios -> Subtly unsafe PII send -> HARD BLOCK")
    pdf.b("Scenarios -> Compounding approve spend -> Hold -> Ops HITL resolve")
    pdf.b("Scenarios -> Policy variant EU refund -> compare US vs EU packs")

    pdf.section("9. Intervention ladder")
    pdf.b("Near-zero -> Pass instantly")
    pdf.b("Low -> Attach receipt")
    pdf.b("Medium -> Soft rewrite")
    pdf.b("High -> Hold at gate (HITL)")
    pdf.b("Critical -> Hard block")

    pdf.section("10. Troubleshooting")
    pdf.b("node/npm not found: install Node 20+ and restart terminal")
    pdf.b("Port 3000 busy: npx next dev -p 3001")
    pdf.b("Old UI: hard refresh Ctrl+Shift+R and restart npm run dev")
    pdf.b("HITL empty: run a high-regret scenario first")

    pdf.section("11. Team / license")
    pdf.p("RegretGate - Accenture Innovation Challenge Round 2 - ControlPlane.ai track")
    pdf.p("Prototype code provided for challenge evaluation.")

    pdf.output(str(path))


def build_proposal_pdf(path: Path) -> None:
    pdf = Doc("RegretGate Business Proposal")
    pdf.title_block(
        "Detailed Business Proposal",
        "RegretGate - Action-aware pre-commit AI control plane  |  ControlPlane.ai",
    )

    pdf.section("1. Problem framing")
    pdf.p(
        "Generative AI already moves money, data, and people inside enterprises - "
        "customer chat, employee copilots, and regulated decision-support. Failures "
        "often surface only after the model has acted."
    )
    pdf.b("Confidently wrong - hallucinations and poor decisions delivered with certainty")
    pdf.b("Quietly expensive - token burn, tool thrash, runaway agent loops")
    pdf.b("Subtly unsafe - bias, leakage, policy violations")
    pdf.p(
        "Today's oversight is fragmented across quality evaluations, cost dashboards, "
        "and safety filters. Problems are detected after commit, when remediation is "
        "costly and liability is real."
    )

    pdf.section("Real-world complexities we design for")
    pdf.b("Different use cases have different latency and risk budgets")
    pdf.b("Bias, hallucination, and privacy risks often overlap")
    pdf.b("No reliable real-time ground truth for every claim")
    pdf.b("Over-flagging creates fatigue; under-flagging creates liability")
    pdf.b("Multi-turn agents compound risk across downstream decisions")
    pdf.b("Regulation differs by geography and industry")
    pdf.b("Foundation models consumed via API - operate at input/output layer")

    pdf.section("Reference operating assumptions")
    pdf.b("Multiple concurrent use cases (support, internal copilot, decision support)")
    pdf.b("Tens of thousands of interactions per week combined")
    pdf.b("Mix of well-governed and loosely-governed knowledge sources")

    pdf.section("2. Solution design")
    pdf.code("Expected Regret = P(failure) x Impact x Irreversibility")
    pdf.p("Core mechanism:")
    pdf.b("Intent and action identification (reply / send / approve / refund / execute / tool loop)")
    pdf.b("Parallel signal analysis: performance, cost/thrash, responsibility")
    pdf.b("Regret estimation: calibrated 0-100 score, use-case and policy aware")
    pdf.b("Responsibility override: hard block for leakage / safety / policy")
    pdf.b("Intervention ladder: pass -> receipt -> soft rewrite -> hold/HITL -> hard block")
    pdf.b("Receipts and audit trail for every decision")
    pdf.b("Feedback loop: human outcomes recalibrate thresholds")

    pdf.section("Differentiation")
    pdf.b("Typical checkers treat every output similarly; RegretGate is action-aware and regret-priced")
    pdf.b("Typical tools are post-hoc; RegretGate is a pre-commit control plane")
    pdf.b("Fixed rules age quickly; RegretGate uses policy packs + human feedback")
    pdf.b("Quality / cost / safety are usually siloed; RegretGate unifies them")

    pdf.section("3. Target users")
    pdf.b("AI Platform / MLOps lead - consistent gate without killing latency")
    pdf.b("Risk / Compliance / Privacy - audit trail, geo packs, hard blocks")
    pdf.b("Business process owners - HITL only where irreversibility is high")
    pdf.b("CIO / Responsible AI sponsors - trust metrics and tunable tradeoffs")

    pdf.section("4. Business case and impact")
    pdf.b("Latency preserved where safe (fast path for low regret)")
    pdf.b("Liability reduced where irreversible (holds and hard blocks)")
    pdf.b("Cost waste reduced via thrash detection")
    pdf.b("Alert fatigue managed by allocating verification proportional to Expected Regret")
    pdf.p(
        "Illustrative framing: 50,000 AI actions/week, ~3% high-regret, ~0.5% "
        "responsibility-critical. Even a small reduction in leakage events or erroneous "
        "automated payouts typically dominates checker operating cost."
    )

    pdf.section("5. Phased roadmap")
    pdf.b("Phase 0 - Prototype (this submission): working PoC, scenarios, ops, docs")
    pdf.b("Phase 1 - Pilot: shadow then enforce, real audit store, PII service, FP/FN labels")
    pdf.b("Phase 2 - Platform: policy studio, compounding graph, thrash budgets, GRC export")
    pdf.b("Phase 3 - Enterprise: multi-tenant controls, SDKs, continuous calibration")

    pdf.section("6. Key risks and mitigations")
    pdf.b("Over-flagging -> ladder + tunable thresholds + fatigue metrics")
    pdf.b("Under-flagging -> responsibility hard block independent of score")
    pdf.b("No ground truth -> receipts, soft rewrite hedges, HITL holds")
    pdf.b("Detector drift / aging rules -> feedback recalibration + versioned policy packs")
    pdf.b("Latency regression -> parallel checks and fast path for low regret")
    pdf.b("Bypass of gate -> platform embedding + audit incompleteness alerts")

    pdf.section("7. What Round 2 demonstrates")
    pdf.p(
        "The prototype proves the core mechanism on sample data: score -> ladder -> "
        "hard block -> HITL -> feedback, across three use cases and three silent failure "
        "modes - without claiming production ML accuracy."
    )
    pdf.p("GitHub: https://github.com/harshalDharpure/controlplane-regretgate")
    pdf.p("Run: npm install && npm run dev  ->  http://localhost:3000")

    pdf.output(str(path))


def set_run(run, size=14, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(5.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*INK)
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], 32, True)

    box2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.5), Inches(1.8))
    tf = box2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = subtitle
    set_run(p2.runs[0], 15, False, MUTED)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*INK)
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.55))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    set_run(tp.runs[0], 22, True)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        if p.runs:
            set_run(p.runs[0], 14, False)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(
        prs,
        "RegretGate",
        "Action-aware AI Control Plane\n"
        "Accenture Innovation Challenge Round 2 - ControlPlane.ai\n"
        "Act with confidence. Verify only when the regret is high.",
    )
    add_bullets_slide(
        prs,
        "1. Problem framing",
        [
            "AI already moves money, data, and people - failures found too late",
            "Confidently wrong · Quietly expensive · Subtly unsafe",
            "Oversight fragmented across quality, cost, and safety tools",
            "Gap: detection after commit, when liability and cost are real",
        ],
    )
    add_bullets_slide(
        prs,
        "Real-world complexities",
        [
            "Different use cases => different latency and risk budgets",
            "Bias, hallucination, and privacy often overlap",
            "No reliable real-time ground truth for every claim",
            "Over-flag fatigue vs under-flag liability must be tuned",
            "Multi-turn agents compound risk; regulation varies by region",
            "Models consumed via API => operate at the input/output layer",
        ],
    )
    add_bullets_slide(
        prs,
        "2. Solution - RegretGate",
        [
            "Expected Regret = P(failure) x Impact x Irreversibility (0-100)",
            "Pre-commit gate on pending actions (reply/send/refund/approve/execute)",
            "Parallel checks: performance, cost/thrash, responsibility",
            "Responsibility risks hard-block (override the normal ladder)",
            "Ladder: Pass -> Receipt -> Soft rewrite -> Hold/HITL -> Hard block",
            "Feedback loop recalibrates thresholds from human decisions",
        ],
    )
    add_bullets_slide(
        prs,
        "Intervention ladder",
        [
            "Near-zero -> Pass instantly (minimal latency)",
            "Low -> Attach receipt (auditability)",
            "Medium -> Soft rewrite (light verification)",
            "High -> Hold at gate (proof or human escalation)",
            "Critical -> Hard block (PII / secrets / unsafe / policy)",
            "Lightweight triage first; deeper verification only when regret is high",
        ],
    )
    add_bullets_slide(
        prs,
        "3. Target users",
        [
            "AI Platform / MLOps - consistent gate without killing latency",
            "Risk / Compliance / Privacy - audit trail, geo packs, hard blocks",
            "Business process owners - HITL only where irreversibility is high",
            "CIO / Responsible AI sponsors - FP/FN proxies and trust metrics",
        ],
    )
    add_bullets_slide(
        prs,
        "4. Business case and impact",
        [
            "Faster where safe; strict where irreversible",
            "Reduce leakage and bad automated payouts before production",
            "Cut quietly expensive agent thrash before spend compounds",
            "Tune over/under-flagging via thresholds + feedback",
            "Illustrative scale: tens of thousands of AI interactions / week",
        ],
    )
    add_bullets_slide(
        prs,
        "5. Phased roadmap",
        [
            "Phase 0 - Prototype (this submission): working PoC + docs",
            "Phase 1 - Pilot: shadow->enforce, audit store, PII, FP/FN labels",
            "Phase 2 - Platform: policy studio, compounding graph, GRC export",
            "Phase 3 - Enterprise: multi-tenant, SDKs, continuous calibration",
        ],
    )
    add_bullets_slide(
        prs,
        "6. Risks and mitigations",
        [
            "Over-flagging -> ladder + tunable thresholds + fatigue metrics",
            "Under-flagging -> responsibility hard block independent of score",
            "No ground truth -> receipts, hedges, HITL holds",
            "Aging rules -> feedback recalibration + versioned policy packs",
            "Latency risk -> parallel checks + fast path for low regret",
        ],
    )
    add_bullets_slide(
        prs,
        "Prototype and how to run",
        [
            "GitHub: github.com/harshalDharpure/controlplane-regretgate",
            "Surfaces: Control Plane (/) · Scenarios (/scenarios) · Ops (/ops)",
            "Run: npm install -> npm run dev -> http://localhost:3000",
            "No API keys required for the core demo",
            "Demonstrates: score -> ladder -> hard block -> HITL -> feedback",
        ],
    )
    add_title_slide(
        prs,
        "Thank you",
        "RegretGate - Minimize Expected Regret\n"
        "Every action scored · Every high-regret action verified\n"
        "Questions welcome",
    )
    prs.save(str(path))


def main():
    readme_pdf = OUT / "RegretGate_README.pdf"
    proposal_pdf = OUT / "RegretGate_Business_Proposal.pdf"
    proposal_pptx = OUT / "RegretGate_Business_Proposal.pptx"

    build_readme_pdf(readme_pdf)
    build_proposal_pdf(proposal_pdf)
    build_pptx(proposal_pptx)

    checklist = OUT / "SUBMISSION_CHECKLIST.txt"
    checklist.write_text(
        "\n".join(
            [
                "Accenture Round 2 - Submission checklist",
                "========================================",
                "",
                "[x] Public GitHub link:",
                "    https://github.com/harshalDharpure/controlplane-regretgate",
                "",
                "[x] README document (PDF):",
                "    RegretGate_README.pdf",
                "",
                "[x] Detailed Business Proposal (PDF):",
                "    RegretGate_Business_Proposal.pdf",
                "",
                "[x] Detailed Business Proposal (PPTX):",
                "    RegretGate_Business_Proposal.pptx",
                "",
                "[ ] Prototype video (mp4 / mov) - RECORD THIS",
                "    Follow docs/DEMO_VIDEO_NOTES.md while running: npm run dev",
                "    Then upload the video on the submission form",
                "",
                "All generated files are in the submission/ folder.",
                "",
            ]
        ),
        encoding="utf-8",
    )

    print("Generated:")
    for p in [readme_pdf, proposal_pdf, proposal_pptx, checklist]:
        print(f" - {p.name}: {p.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
